from pptx import Presentation
from pptx.util import Inches, Px
from pptx.constants import MSO

templateName = 'template_red.pptx'

prs = Presentation(templateName)
#title_slidelayout = prs.slidemasters[0].slidelayouts[0]
#slide = prs.slides.add_slide(title_slidelayout)
#title = slide.shapes.title
#subtitle = slide.shapes.placeholders[1]

#title.text = "Tittle"
#subtitle.text = "subtitle"

#------------------------------------------
slidelayout = []
nMaster = 0
for j in range(0, 20):
    try:
        prs.slidemasters[j]
    except IndexError:
        print "only %d masterSlideLayouts" % (j)
        nMaster = j
        break
for j in range(0, nMaster):
    for i in range(0, 25):
        try:
            slidelayout.append(prs.slidemasters[j].slidelayouts[i])
            slide = prs.slides.add_slide(slidelayout[i])
            shapes = slide.shapes
            shape = shapes.add_shape(MSO.CHART, 5,5,5,5)
            try:
                title = slide.shapes.title
                title.text = "masterSlide %d: slideLayout %d" % (j, i)
            except AttributeError:
                print "slideLayout %d does not have title placeholder" % i
        except IndexError:
            print "masterSlide %d only has %d slideLayouts" % (j, i)
            break


#------------------------------------------
#img_path = '011-Shear-x=0.jpg'

#blank_slidelayout = prs.slidemasters[1].slidelayouts[0]
#slide = prs.slides.add_slide(blank_slidelayout)

#left = top = Inches(5)
#pic = slide.shapes.add_picture(img_path, left, top)

#left = Inches(5)
#width = Px(120)
#height = int(width * 1.427)
#pic = slide.shapes.add_picture(img_path, left, top, width)

#--------------------------------------------

prs.save('%s-layouts.pptx' % templateName.split('.')[0])
