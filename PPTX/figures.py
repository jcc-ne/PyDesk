from pptx import Presentation
from pptx.util import Cm, Pt
import os
import glob
from parseJpg import *
# from pptx.util import Inches, Px


def main():
    prs = Presentation('template_red.pptx')
    title_slidelayout = prs.slidemasters[0].slidelayouts[0]
    slide = prs.slides.add_slide(title_slidelayout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]

    title.text = "Title!"
    subtitle.text = "subtitle"

#------------------------------------------
    dirs = [o for o in glob.glob('../../Run/*') if os.path.isdir(o)]
    print "dirs=", dirs

    for dir in dirs:
        print "\n ------ in dir ", dir, "------"
        psJpgs = searchPsJpgs(dir)
        frames = sortPsJpgs(psJpgs)
        for key in frames:
            print frames[key]
            if len(frames[key]) == 2:
                addTwoFigs(prs, frames[key][0],
                           frames[key][1])

    prs.save('printout_cfdresults.pptx')


def addTwoFigs(prs, psjpg1, psjpg2, titleText=""):
    slidelayout = prs.slidemasters[1].slidelayouts[0]
    slide = prs.slides.add_slide(slidelayout)

    #---- add figure on the left --------------
    img_path = psjpg1.jpgfileFp

    left = Cm(1.27)
    top = Cm(5.15)
    width = Cm(13.27)
    # height = Cm(19.05)
    slide.shapes.add_picture(img_path, left, top, width)

    #---- add figure on the right--------------
    img_path = psjpg2.jpgfileFp

    left = Cm(10.86)
    slide.shapes.add_picture(img_path, left, top, width)

    #---- title ------
    title = slide.shapes.title
    if titleText == "":
        titleText = psjpg1.field
    title.text = titleText + " Field"
    addTextBot(slide, "%s at plane %s" % (psjpg1.fieldFull, psjpg1.loc))
    addTable(slide)
    return slide


def addThreeFigs(prs, titleText, img_path_1, img_path_2, img_path_3):
    slidelayout = prs.slidemasters[1].slidelayouts[1]
    slide = prs.slides.add_slide(slidelayout)
    title = slide.shapes.title
    title.text = titleText

    #---- add figure on the left --------------
    img_path = img_path_1

    left = Cm(1.27)
    top = Cm(5.15)
    width = Cm(13.27)
    # height = Cm(19.05)
    slide.shapes.add_picture(img_path, left, top, width)

    #---- add figure in the middle--------------
    img_path = img_path_2

    left = Cm(10.86)
    slide.shapes.add_picture(img_path, left, top, width)

    #---- add figure on the very right--------------
    img_path = img_path_3

    left = Cm(14.86)
    slide.shapes.add_picture(img_path, left, top, width)


def addTable(slide):
    shapes = slide.shapes
    rows = 1
    cols = 2
    left = Cm(1.27)
    top = Cm(10)
    width = Cm(6.0)
    height = Cm(0.8)

    tbl = shapes.add_table(rows, cols, left, top, width, height)

    # set column widths
    tbl.columns[0].width = Cm(9.59)
    tbl.columns[1].width = Cm(13.27)

    # write column headings
    tbl.cell(0, 0).text = 'Foo'
    tbl.cell(0, 1).text = 'Bar'

    # write body cells
    # tbl.cell(1, 0).text = 'Baz'
    # tbl.cell(1, 1).text = 'Qux'


def addTextBot(slide, boxText):
    #---- add text box on the bottom--------------
    left = Cm(0.0)
    top = Cm(17.55)
    width = Cm(25.4)
    height = Cm(1.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.textframe
    tf.text = boxText
    tf.paragraphs[0].font.size = Pt(16)
    tf.bold = False

    #p = tf.add_paragraph()
    #p.text = boxText
    #p.font.bold = False

    #p = tf.add_paragraph()
    #p.text = "This is a third paragraph that's big"
    #p.font.size = Pt(40)
    #f = txBox.textframe


if __name__ == "__main__":
    main()
