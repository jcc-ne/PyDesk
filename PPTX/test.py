from pptx import Presentation

prs = Presentation()
title_slidelayout = prs.slidelayouts[0]
slide = prs.slides.add_slide(title_slidelayout)
title = slide.shapes.title
subtitle = slide.shapes.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

#------------------------------------------

from pptx import Presentation
from pptx.util import Inches, Px

img_path = 'monty-truth.png'

prs = Presentation()
blank_slidelayout = prs.slidelayouts[6]
slide = prs.slides.add_slide(blank_slidelayout)

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(5)
width = Px(280)
height = int(width*1.427)
pic = slide.shapes.add_picture(img_path, left, top, width, height)

#--------------------------------------------

prs.save('test.pptx')
