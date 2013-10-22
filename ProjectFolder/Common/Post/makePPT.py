from pptx import Presentation
from pptx.oxml import _SubElement
from pptx.util import Cm, Pt
import os
import glob
from parseJpg import searchAllPsJpgs, sortGrpPsJpgs
import Image


def main():
    prs = Presentation('template_red.pptx')
    title_slidelayout = prs.slidemasters[0].slidelayouts[0]
    slide = prs.slides.add_slide(title_slidelayout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]

    title.text = "Title!"
    subtitle.text = "subtitle"

    #-------glob current folder
    Dirs = ['./fig/']
    psJpgs = searchAllPsJpgs(Dirs)
    for psJpg in psJpgs:
        psJpg.printAll()
    allSlides = sortGrpPsJpgs(psJpgs)
  #  slidesEachNumField(prs, allSlides)
    slidesCompareFields(prs, allSlides)
    slidesCompareNum(prs, allSlides)

    #------------------------------------------
    Dirs = [o for o in glob.glob('../../Run/*') if os.path.isdir(o)]
    for Dir in Dirs:
        Dir = Dir.replace("\\", "/")
    print "Dirs=", Dirs
    psJpgs = searchAllPsJpgs(Dirs)
    allSlides = sortGrpPsJpgs(psJpgs)

    slidesEachNumField(prs, allSlides)
    slidesCompareFields(prs, allSlides)
    slidesCompareNum(prs, allSlides)

    foutName = 'printout_cfdresults.pptx'
    prs.save(foutName)


def slidesEachNumField(prs, allSlides):
    priors = []
    groupKeys = []
    priors.append('case')
    priors.append('num')
    groupKeys.append('case')    # things they have in common
    groupKeys.append('num')     # things they have in common
    groupKeys.append('field')   # things they have in common
    titleKeys = ['case', 'numFull', 'fieldFull']     # keys to determine
                                                     # the title of slide

    tabKeys = ['locFull']  # will determine the text in the table

    slides = allSlides.sortWithNewKeys(priors, groupKeys,
                                       titleKeys, tabKeys)
    countFigsNmakeSlides(prs, slides)


def slidesCompareFields(prs, allSlides):
    priors = []
    groupKeys = []
    priors.append('case')
    priors.append('num')
    groupKeys.append('case')    # things they have in common
    groupKeys.append('num')     # things they have in common
    titleKey = 'numFull'             # will determine the title of slide
    titleKeys = ['case', 'numFull']     # keys to determine the title of slide
    groupKeys.append(titleKey.replace('Full', ''))

    tabKeys = ['fieldFull', 'unit']  # will determine the text in the table

    slides = allSlides.sortWithNewKeys(priors, groupKeys, titleKeys, tabKeys)
    countFigsNmakeSlides(prs, slides)


def slidesCompareNum(prs, allSlides):
    priors = []
    groupKeys = []
    priors.append('case')
    priors.append('field')
    groupKeys.append('case')    # things they have in common
    groupKeys.append('field')     # things they have in common
    groupKeys.append('loc')     # things they have in common
    titleKey = 'fieldFull'             # will determine the title of slide
    titleKeys = ['case', 'fieldFull']     # keys to determine
                                          # the title of slide
    groupKeys.append(titleKey.replace('Full', ''))
    tabKeys = ['numFull']  # will determine the text in the table

    slides = allSlides.sortWithNewKeys(priors, groupKeys, titleKeys, tabKeys)
    countFigsNmakeSlides(prs, slides)


def countFigsNmakeSlides(prs, slidepages):
    i = 0
    for s in slidepages:
        i += 1
        print "\n +++ slidepage %d: " % i,
        print [s.frames[j].code for j in range(s.nf)]
        if s.nf == 2:
            addTwoFigs(prs, s)

        if s.nf == 6:
            addSixFigs(prs, s)


def imageSize(img_path):
    img = Image.open(img_path)
    pixelWidth = float(img.size[0])
    pixelHeight = float(img.size[1])
    # dpi = float(img.info['dpi'])
    # cmWidth = Cm(pixelWidth / dpi * 2.54)
    # cmHeight = Cm(pixelHeight / dpi * 2.54)
    return [pixelWidth, pixelHeight]


def addTwoFigs(prs, slidepage, titleText=""):
    slidelayout = prs.slidemasters[1].slidelayouts[0]
    slide = prs.slides.add_slide(slidelayout)

    #---- add figure on the left --------------
    img_path = slidepage.frames[0].jpgfileFp

    top = Cm(6.25)
    tabLeft = left = Cm(1.14)
    width = Cm(11.56)
    # height = Cm(19.05)
    slide.shapes.add_picture(img_path, left, top, width)

    #---- add figure on the right--------------
    img_path = slidepage.frames[1].jpgfileFp

    left = left + width
    slide.shapes.add_picture(img_path, left, top, width)

    #---- title ------
    title = slide.shapes.title
    title.text = slidepage.titleText
    print "     title = ", titleText
    addTableTwo(slide, tabLeft,
                width, slidepage.tabTextList)
    addTextBot(slide, slidepage.boxText, left=tabLeft)
    return slide


def addThreeFigs(prs, psjpg1, psjpg2, psjpg3, titleText=""):
    slidelayout = prs.slidemasters[1].slidelayouts[0]
    slide = prs.slides.add_slide(slidelayout)

    #---- add figure on the left --------------
    img_path = psjpg1.jpgfileFp

    txBoxLeft = left = Cm(1.27)
    top = Cm(4.65)
    width = Cm(13.27)
    # height = Cm(19.05)
    slide.shapes.add_picture(img_path, left, top, width)

    #---- add figure in the middle--------------
    img_path = psjpg2.jpgfileFp

    left = Cm(8.86)
    slide.shapes.add_picture(img_path, left, top, width)

    #---- add figure on the right--------------
    img_path = psjpg3.jpgfileFp

    left = Cm(12.86)
    slide.shapes.add_picture(img_path, left, top, width)

    #---- title ------
    title = slide.shapes.title
    if titleText == "":
        titleText = psjpg1.case + ': ' + psjpg1.fieldFull
        title.text = titleText
        addTextBot(slide, "Operating point: %s" %
                   (psjpg1.numFull), left=txBoxLeft)
        addTableTwo(slide, psjpg1, psjpg2)
        return slide


def addSixFigs(prs, slidepage, titleText=""):
    slidelayout = prs.slidemasters[1].slidelayouts[0]
    slide = prs.slides.add_slide(slidelayout)
    #width = Cm(8.03)
    height = Cm(6)

    #---- add figure on the left --------------
    img_path = slidepage.frames[0].jpgfileFp
    [picW, picH] = imageSize(img_path)
    WHratio = picW / picH
    width = height * WHratio
    top = Cm(4.65)

    txBoxLeft = left = Cm(1.27)
    left = Cm(2.5)
    # height = Cm(19.05)
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)

    #---- add figure in the middle--------------
    img_path = slidepage.frames[1].jpgfileFp

    left = Cm(9.25)
    slide.shapes.add_picture(img_path, left, top, height=height)

    #---- add figure on the right--------------
    img_path = slidepage.frames[2].jpgfileFp

    left = Cm(16)
    slide.shapes.add_picture(img_path, left, top, height=height)

    #---- add figure on the leftBot --------------
    img_path = slidepage.frames[3].jpgfileFp
    top = Cm(11.33)

    left = Cm(2.5)
    slide.shapes.add_picture(img_path, left, top, height=height)

    #---- add figure on the midBot--------------
    img_path = slidepage.frames[4].jpgfileFp
    left = Cm(9.25)

    slide.shapes.add_picture(img_path, left, top, height=height)

    #---- add figure on the rightBot--------------
    img_path = slidepage.frames[5].jpgfileFp
    left = Cm(16)

    slide.shapes.add_picture(img_path, left, top, height=height)

    #---- title ------
    title = slide.shapes.title
    title.text = slidepage.titleText
    print "     title = ", titleText
    # addTextBot(slide, "Operating point: %s" % (psjpg1.op), left=txBoxLeft)
    addTableSix(slide,
                width, slidepage.tabTextList)

    return slide


def addTableTwo(slide, left, tabwidth, tabTextList):
    shapes = slide.shapes
    rows = 1
    cols = 2
    # left = Cm(1.27)
    top = Cm(16.4)
    tabwidth = int(tabwidth)  # pixel has to be integer
    # width = Cm(22.88)
    width = tabwidth * 2
    height = Cm(0.8)

    tbl = shapes.add_table(rows, cols, left, top, width, height)

    # set column widths
    tbl.columns[0].width = tabwidth
    tbl.columns[1].width = tabwidth

    for i in range(0, 2):
        text = tabTextList[i]
        print "          tabText = %s" % text
        tbl.cell(0, i).text = text
        tf = tbl.cell(0, i).textframe
        font = tf.paragraphs[0].font
        font.size = Pt(16)
        font.bold = False
        set_font_color_and_typeface(font, rgbColor("BLACK"), 'Arial')


def addTableSix(slide,
                tabwidth,
                tabTextList):
    shapes = slide.shapes
    rows = 1
    cols = 3
    left = Cm(2.5)
    top = Cm(10.36)
    tabwidth = int(tabwidth)  # pixel has to be integer
    width = tabwidth * 3
    height = Cm(0.6)

    tbl = shapes.add_table(rows, cols, left, top, width, height)

    # set column widths
    tbl.columns[0].width = tabwidth
    tbl.columns[1].width = tabwidth
    tbl.columns[2].width = width - 2 * tabwidth

    # write column headings

    for i in range(0, 3):
        text = tabTextList[i]
        print "          tabText = %s" % text
        tbl.cell(0, i).text = text
        tf = tbl.cell(0, i).textframe
        font = tf.paragraphs[0].font
        font.size = Pt(10)
        font.bold = False
        set_font_color_and_typeface(font, rgbColor("BLACK"), 'Arial')

    top = Cm(16.93)
    tbl = shapes.add_table(rows, cols, left, top, width, height)

    for i in range(0, 3):
        text = tabTextList[i + 2]
        print "          tabText = %s" % text
        tbl.cell(0, i).text = text
        tf = tbl.cell(0, i).textframe
        font = tf.paragraphs[0].font
        font.size = Pt(10)
        font.bold = False
        set_font_color_and_typeface(font, rgbColor("BLACK"), 'Arial')


def addTextBot(slide, boxText, left=Cm(0.0), top=Cm(17.55)):
    #---- add text box on the bottom--------------
    # left = Cm(0.0)
    # top = Cm(17.55)
    width = Cm(25.4)
    height = Cm(1.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.textframe
    tf.text = boxText
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = False

    #p = tf.add_paragraph()
    #p.text = boxText
    #p.font.bold = False

    #p = tf.add_paragraph()
    #p.text = "This is a third paragraph that's big"
    #p.font.size = Pt(40)
    #f = txBox.textframe


def rgbColor(colorName):
    colorlist = {}
    colorlist["ORANGE"] = 'FF6600'
    colorlist["WHITE"] = 'FFFFFF'
    colorlist["BLACK"] = '000000'
    return colorlist[colorName]


def set_font_color_and_typeface(font, rgbColor, typeface=None):
    rPr = font._Font__rPr
    solidFill = _SubElement(rPr, 'a:solidFill')
    srgbClr = _SubElement(solidFill, 'a:srgbClr')
    srgbClr.set('val', rgbColor)
    if typeface:
        latin = _SubElement(rPr, 'a:latin')
        latin.set('typeface', typeface)


if __name__ == "__main__":
    main()
